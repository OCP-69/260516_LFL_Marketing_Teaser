#!/usr/bin/env python3
"""Translate LFL_Kostentransparenz_Interpack.pptx to English in-place copy."""
from __future__ import annotations

import csv
import json
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

BASE = Path(__file__).resolve().parent
SRC = BASE / "LFL_Kostentransparenz_Interpack.pptx"
OUT_PPTX = BASE / "LFL_Kostentransparenz_Interpack_EN.pptx"
OUT_MD = BASE / "LFL_Kostentransparenz_Interpack_EN_Translation.md"
OUT_CSV = BASE / "LFL_Kostentransparenz_Interpack_EN_Translation.csv"

# Exact source strings from slide (incl. PowerPoint line breaks \u000b)
TRANSLATIONS: dict[str, str] = {
    "LFL": "LFL",
    "LoopForgeLab": "LoopForgeLab",
    "Forge Engine · Cost Intelligence": "Forge Engine · Cost Intelligence",
    "Teaser · Interpack\u000bInsights Verpackungs- & Sondermaschinenbau": (
        "Teaser · Interpack\u000bInsights — packaging & custom machinery"
    ),
    "Forge Engine · Kosten-Intelligence aus Ihren Daten": (
        "Forge Engine · cost intelligence from your data"
    ),
    "80 % der Kosten im Design — bewertet wird erst Wochen später.": (
        "80% of costs are set in design — yet evaluated weeks later."
    ),
    "Belastbare Kosten aus CAD, BOM und Einkauf — ohne Migration. Forge Engine liefert Stückkosten und TCO im Designmoment, in Stunden statt Wochen, ohne ERP-Projekt.": (
        "Reliable costs from CAD, BOM and purchasing — no migration. "
        "Forge Engine delivers unit costs and TCO at design stage, in hours not weeks, "
        "without an ERP project."
    ),
    "14–28  →  5": "14–28  →  5",
    "Tage / RFQ": "days / RFQ",
    "< € 2.500": "< € 2,500",
    "pro Kalkulationszyklus": "per costing cycle",
    "3D- & CAD-Daten als Ausgangspunkt": "3D & CAD data as the starting point",
    "Der Wechsel": "The shift",
    "Datenschatz ungenutzt  →  80 % alter BOMs sofort nutzbar": (
        "Untapped data  →  80% of legacy BOMs usable immediately"
    ),
    "80 %": "80%",
    "der Lebenszyklus-\u000bkosten im Design": "of lifecycle\u000bcosts at design stage",
    "WAS SIE GEWINNEN": "WHAT YOU GAIN",
    "EINSTIEG IN MEHR KLARHEIT: IHR PILOTPROJEKT": (
        "GET STARTED WITH MORE CLARITY: YOUR PILOT PROJECT"
    ),
    "Marge": "Margin",
    "Tempo": "Speed",
    "Datenschatz": "Data assets",
    "Sicherheit": "Certainty",
    "Puffer durch belastbare Zahlen ersetzen — schärfere Angebote ohne Marge zu verschenken.": (
        "Replace buffers with reliable numbers — sharper quotes without giving away margin."
    ),
    "Varianten in Stunden statt Wochen. RFQ 14–28 → 5 Tage — mehr Anfragen ohne mehr Köpfe.": (
        "Variants in hours not weeks. RFQ 14–28 → 5 days — more inquiries without more headcount."
    ),
    "Alte Stücklisten und Einkaufspreise nutzbar — 80 % Wiederverwendung statt Kalkulation von Null.": (
        "Legacy BOMs and purchase prices usable — 80% reuse instead of costing from scratch."
    ),
    "Material, Lieferant, Regulierung pro Markt — sofort durchgerechnet, nicht in der Krise.": (
        "Material, supplier, regulation per market — costed instantly, not in a crisis."
    ),
    "Intelligenzschicht auf bestehende Systeme. Kein ERP-Replace, kein IT-Rollout. Erste Design-Partner Q2 2026.": (
        "Intelligence layer on existing systems. No ERP replacement, no IT rollout. "
        "First design partners Q2 2026."
    ),
    "On-Prem möglich": "On-prem available",
    "LFL: Team-Erfahrung:  KI · SaaS · Industrie\u000bLösung relevant für: Geschäftsführung · Design/Engineering · Einkauf · Controller · Vertrieb": (
        "LFL team experience: AI · SaaS · industry\u000b"
        "Relevant for: executive management · design/engineering · purchasing · controlling · sales"
    ),
    "Beispiel: Kostenanteile je Dimension — Material · Prozess · Rohstoff": (
        "Example: cost share by dimension — material · process · raw material"
    ),
    "Erste Design-Partner 2026 · begrenzte Plätze": (
        "First design partners 2026 · limited slots"
    ),
    "Weniger Iteration. Mehr Klarheit \u000b— bevor das nächste Angebot an fehlenden Informationen scheitert.": (
        "Less iteration. More clarity\u000b— before the next quote fails on missing information."
    ),
    "LoopForgeLab· Cost Intelligence  ·  hello@loopforgelab.com · Olaf Pick · Berlin": (
        "LoopForgeLab · Cost Intelligence · hello@loopforgelab.com · Olaf Pick · Berlin · Germany"
    ),
    "DSGVO-konform": "GDPR-compliant",
    "EU-Datenhaltung": "EU data hosting",
    "3–10 BOMs + STEP sind ausreichend zum Start": (
        "3–10 BOMs + STEP files are enough to get started"
    ),
    "Ihr Ansprechpartner": "Your contact",
    "Olaf Pick\u000bCo-Founder LoopForgeLab": "Olaf Pick\u000bCo-Founder LoopForgeLab",
    "olaf@loopforgelab.com": "olaf@loopforgelab.com",
    "Erfahren Sie, wie unsere Lösung Ihre Kosten und Ihren Aufwand reduzieren hilft": (
        "See how our solution helps reduce your costs and effort"
    ),
    "Terminbuchung\u000b15-Min-Gespräch \u000bcalendly.com\\olaf-pick\\olaf-15min": (
        "Book a meeting\u000b15-min call\u000bcalendly.com/olaf-pick/olaf-15min"
    ),
}


def walk_shapes(shapes):
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from walk_shapes(shape.shapes)
        elif hasattr(shape, "text_frame"):
            yield shape
        elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            for row in shape.table.rows:
                for cell in row.cells:
                    yield cell


def set_text(shape, new_text: str) -> None:
    tf = shape.text_frame
    lines = new_text.split("\u000b")
    tf.paragraphs[0].text = lines[0]
    for line in lines[1:]:
        tf.add_paragraph().text = line


def main():
    prs = Presentation(str(SRC))
    rows: list[dict] = []
    missing: list[str] = []

    for slide in prs.slides:
        for shape in walk_shapes(slide.shapes):
            if not hasattr(shape, "text_frame"):
                continue
            original = shape.text
            if not original.strip():
                continue
            if original in TRANSLATIONS:
                english = TRANSLATIONS[original]
                set_text(shape, english)
                rows.append({"original": original, "english": english})
            else:
                missing.append(original)

    prs.save(str(OUT_PPTX))

    # Markdown table
    def esc(s: str) -> str:
        return s.replace("\u000b", "<br>").replace("|", "\\|").replace("\n", "<br>")

    md_lines = [
        "# LFL Kostentransparenz Interpack — Übersetzung EN",
        "",
        f"Quelle: `{SRC.name}` → `{OUT_PPTX.name}`",
        "",
        "| Original (DE) | English |",
        "|---|---|",
    ]
    for r in rows:
        md_lines.append(f"| {esc(r['original'])} | {esc(r['english'])} |")
    if missing:
        md_lines.extend(["", "## Nicht zugeordnet", ""])
        for m in missing:
            md_lines.append(f"- {esc(m)}")
    OUT_MD.write_text("\n".join(md_lines), encoding="utf-8")

    with OUT_CSV.open("w", encoding="utf-8-sig", newline="") as f:
        w = csv.DictWriter(f, fieldnames=["original", "english"])
        w.writeheader()
        w.writerows(rows)

    print(f"PPTX: {OUT_PPTX}")
    print(f"MD:   {OUT_MD}")
    print(f"CSV:  {OUT_CSV}")
    print(f"Translated: {len(rows)} | Missing: {len(missing)}")
    if missing:
        Path(BASE / "_missing_text.json").write_text(
            json.dumps(missing, ensure_ascii=False, indent=2), encoding="utf-8"
        )


if __name__ == "__main__":
    main()
